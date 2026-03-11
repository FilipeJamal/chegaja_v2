from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


WORKSPACE = Path(r"C:\Users\HP\Documents\ProjetosFlutter\chegaja_v2")
OUT_DIR = WORKSPACE / "artifacts" / "presentation_chegaja"
OUT_FILE = OUT_DIR / "ChegaJa_Apresentacao_App_COMPAT.pptx"
OUT_FILE_FULL = OUT_DIR / "ChegaJa_Apresentacao_App_COMPLETA_COMPAT.pptx"
OUT_FILE_FIXED = OUT_DIR / "ChegaJa_Apresentacao_App_CHAT_COMPLETA_COMPAT.pptx"

IMG_ICON = WORKSPACE / "assets" / "images" / "app_icon.png"

IMG_ROLE = OUT_DIR / "role_selector_open.png"
IMG_CLIENT_HOME = OUT_DIR / "client_home.png"
IMG_CLIENT_ORDERS = OUT_DIR / "client_orders.png"
IMG_CLIENT_MESSAGES = OUT_DIR / "client_messages_open.png"
IMG_CLIENT_CHAT_THREAD = OUT_DIR / "client_chat_thread.png"
IMG_CLIENT_PROFILE = OUT_DIR / "client_profile_fixed.png"
IMG_CLIENT_SUPPORT = OUT_DIR / "client_support.png"
IMG_PROVIDER_HOME = OUT_DIR / "provider_home.png"
IMG_PROVIDER_JOBS = OUT_DIR / "provider_jobs_fixed.png"
IMG_PROVIDER_MESSAGES = OUT_DIR / "provider_messages_fixed.png"
IMG_PROVIDER_CHAT_THREAD = OUT_DIR / "provider_chat_thread.png"
IMG_PROVIDER_PROFILE = OUT_DIR / "provider_profile.png"
IMG_PROVIDER_PAYMENTS = OUT_DIR / "provider_payments.png"
IMG_PROVIDER_SUPPORT = OUT_DIR / "provider_support.png"
IMG_ADMIN = OUT_DIR / "admin_panel_fixed.png"


BG = RGBColor(0xF6, 0xF8, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x11, 0x14, 0x18)
MUTED = RGBColor(0x6B, 0x72, 0x80)
PRIMARY = RGBColor(0x12, 0xBA, 0x9B)
NAVY = RGBColor(0x0B, 0x3C, 0x5D)
BORDER = RGBColor(0xDD, 0xE5, 0xE8)
SOFT = RGBColor(0xEE, 0xF7, 0xF5)
AMBER = RGBColor(0xF3, 0xB4, 0x2F)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_background(slide, color=BG):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text="",
    font_size=18,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    font_name="Arial",
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, font_size=12, color=MUTED):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(2)
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.bullet = True
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, title, body=None, bullets=None, fill=WHITE):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = BORDER
    add_textbox(
        slide,
        x + Inches(0.12),
        y + Inches(0.10),
        w - Inches(0.24),
        Inches(0.26),
        title,
        font_size=16,
        bold=True,
        color=NAVY,
    )
    if body:
        add_textbox(
            slide,
            x + Inches(0.12),
            y + Inches(0.42),
            w - Inches(0.24),
            h - Inches(0.50),
            body,
            font_size=12,
            color=MUTED,
        )
    if bullets:
        add_bullets(
            slide,
            x + Inches(0.08),
            y + Inches(0.36),
            w - Inches(0.16),
            h - Inches(0.42),
            bullets,
        )
    return card


def add_header(slide, eyebrow, title, subtitle=""):
    add_textbox(
        slide,
        Inches(0.55),
        Inches(0.35),
        Inches(2.8),
        Inches(0.3),
        eyebrow.upper(),
        font_size=11,
        bold=True,
        color=PRIMARY,
    )
    add_textbox(
        slide,
        Inches(0.55),
        Inches(0.62),
        Inches(8.9),
        Inches(0.55),
        title,
        font_size=24,
        bold=True,
        color=TEXT,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.55),
            Inches(1.10),
            Inches(11.9),
            Inches(0.55),
            subtitle,
            font_size=12,
            color=MUTED,
        )


def add_picture_fit(slide, image_path, x, y, w, h):
    with Image.open(image_path) as img:
        iw, ih = img.size
    image_ratio = iw / ih
    box_ratio = w / h
    if image_ratio > box_ratio:
        new_w = w
        new_h = w / image_ratio
        new_x = x
        new_y = y + (h - new_h) / 2
    else:
        new_h = h
        new_w = h * image_ratio
        new_x = x + (w - new_w) / 2
        new_y = y
    slide.shapes.add_picture(str(image_path), new_x, new_y, new_w, new_h)


def add_phone_frame(slide, image_path, x, y, w, h, caption, caption_color=MUTED):
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = BORDER
    add_picture_fit(
        slide,
        image_path,
        x + Inches(0.08),
        y + Inches(0.08),
        w - Inches(0.16),
        h - Inches(0.24),
    )
    add_textbox(
        slide,
        x,
        y + h + Inches(0.02),
        w,
        Inches(0.25),
        caption,
        font_size=10,
        color=caption_color,
        align=PP_ALIGN.CENTER,
    )


def add_gallery_2x2(slide, items, start_x=0.75, start_y=1.75, box_w=2.65, box_h=2.2):
    gap_x = 0.28
    gap_y = 0.45
    for index, (image_path, caption) in enumerate(items):
        col = index % 2
        row = index // 2
        x = Inches(start_x + col * (box_w + gap_x))
        y = Inches(start_y + row * (box_h + gap_y))
        add_phone_frame(slide, image_path, x, y, Inches(box_w), Inches(box_h), caption)


def add_flow_bar(slide, x, y, steps):
    gap = 0.16
    width = (11.95 - (len(steps) - 1) * gap) / len(steps)
    for index, (title, fill) in enumerate(steps):
        current_x = x + Inches(index * (width + gap))
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            current_x,
            y,
            Inches(width),
            Inches(0.56),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = fill
        add_textbox(
            slide,
            current_x + Inches(0.05),
            y + Inches(0.10),
            Inches(width - 0.10),
            Inches(0.26),
            title,
            font_size=11,
            bold=True,
            color=WHITE if fill == NAVY else TEXT,
            align=PP_ALIGN.CENTER,
        )


def slide_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(
        slide,
        Inches(0.65),
        Inches(0.65),
        Inches(6.1),
        Inches(0.35),
        "MARKETPLACE DE SERVICOS",
        font_size=11,
        bold=True,
        color=PRIMARY,
    )
    add_textbox(
        slide,
        Inches(0.65),
        Inches(1.05),
        Inches(6.0),
        Inches(1.2),
        "ChegaJa\nApresentacao completa do aplicativo",
        font_size=26,
        bold=True,
        color=TEXT,
    )
    add_textbox(
        slide,
        Inches(0.65),
        Inches(2.32),
        Inches(5.9),
        Inches(0.78),
        "Deck focado no funcionamento do produto, no fluxo ponta a ponta e nas telas reais disponiveis do app.",
        font_size=13,
        color=MUTED,
    )
    add_card(
        slide,
        Inches(0.65),
        Inches(3.08),
        Inches(5.45),
        Inches(2.18),
        "O que entra nesta versao",
        bullets=[
            "visao geral do produto e proposta de valor",
            "fluxo do cliente, mensagens, suporte e perfil",
            "fluxo do prestador, trabalhos, pagamentos e KYC",
            "backoffice, administracao e controles operacionais",
        ],
    )
    add_card(slide, Inches(8.18), Inches(1.00), Inches(4.05), Inches(4.70), "Identidade do produto")
    add_picture_fit(slide, IMG_ICON, Inches(8.58), Inches(1.48), Inches(3.2), Inches(3.2))
    add_textbox(
        slide,
        Inches(8.35),
        Inches(4.96),
        Inches(3.7),
        Inches(0.36),
        "ChegaJa v2",
        font_size=16,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(8.28),
        Inches(5.28),
        Inches(3.82),
        Inches(0.40),
        "Marketplace mobile para servicos locais",
        font_size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def slide_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Visao geral",
        "O que e o ChegaJa",
        "Aplicativo que liga clientes, prestadores e operacao num fluxo unico: descoberta, pedido, comunicacao, execucao, pagamento e suporte.",
    )
    add_card(
        slide,
        Inches(0.65),
        Inches(1.8),
        Inches(2.95),
        Inches(2.1),
        "Objetivo",
        bullets=[
            "reduzir friccao na contratacao",
            "dar previsibilidade ao cliente",
            "gerar demanda ao prestador",
        ],
    )
    add_card(
        slide,
        Inches(3.82),
        Inches(1.8),
        Inches(2.6),
        Inches(2.1),
        "Perfis",
        bullets=["cliente", "prestador", "admin"],
    )
    add_card(
        slide,
        Inches(6.64),
        Inches(1.8),
        Inches(2.85),
        Inches(2.1),
        "Tipos de pedido",
        bullets=["imediato", "agendado", "por orcamento"],
    )
    add_card(
        slide,
        Inches(9.71),
        Inches(1.8),
        Inches(2.95),
        Inches(2.1),
        "Base tecnica",
        bullets=["Flutter", "Firebase", "Stripe", "localizacao em tempo real"],
    )
    add_card(
        slide,
        Inches(0.65),
        Inches(4.2),
        Inches(12.0),
        Inches(2.1),
        "Leitura do produto",
        bullets=[
            "o cliente entra pelo servico e acompanha tudo pelo pedido e pelas mensagens",
            "o prestador gere disponibilidade, trabalhos, conversas, pagamentos e validacao documental",
            "a operacao acompanha indicadores, suporte, no-show e moderacao da plataforma",
        ],
    )


def slide_entry():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Entrada",
        "Primeira decisao do utilizador: que papel vai assumir",
        "A experiencia muda logo no inicio conforme o papel escolhido no ecossistema.",
    )
    add_phone_frame(
        slide,
        IMG_ROLE,
        Inches(0.85),
        Inches(1.7),
        Inches(4.2),
        Inches(4.75),
        "Tela real de selecao de papel",
    )
    add_card(
        slide,
        Inches(5.4),
        Inches(1.82),
        Inches(3.2),
        Inches(1.9),
        "Como funciona",
        bullets=[
            "autenticacao leve para entrada rapida",
            "cliente e prestador seguem jornadas diferentes",
            "tabs e funcionalidades mudam por perfil",
        ],
    )
    add_card(
        slide,
        Inches(8.9),
        Inches(1.82),
        Inches(3.3),
        Inches(1.9),
        "Navegacao principal",
        bullets=[
            "cliente: Inicio, Pedidos, Mensagens, Perfil",
            "prestador: Inicio, Trabalhos, Mensagens, Perfil",
            "admin: painel operacional dedicado",
        ],
    )
    add_card(
        slide,
        Inches(5.4),
        Inches(4.05),
        Inches(6.8),
        Inches(1.65),
        "Importancia desta tela",
        body="O produto nao obriga o utilizador a atravessar um onboarding pesado antes de mostrar valor. Primeiro identifica a intencao; depois abre a experiencia adequada.",
    )


def slide_client_gallery_1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Cliente",
        "Telas principais da jornada do cliente",
        "Aqui o foco e descobrir servicos, consultar pedidos, acompanhar mensagens e gerir a conta.",
    )
    add_gallery_2x2(
        slide,
        [
            (IMG_CLIENT_HOME, "Home do cliente"),
            (IMG_CLIENT_ORDERS, "Pedidos"),
            (IMG_CLIENT_MESSAGES, "Mensagens"),
            (IMG_CLIENT_PROFILE, "Perfil"),
        ],
    )
    add_card(
        slide,
        Inches(6.35),
        Inches(1.75),
        Inches(6.0),
        Inches(4.95),
        "Leitura do fluxo do cliente",
        bullets=[
            "a home expõe categorias e modos de contratacao",
            "a aba de pedidos concentra o historico e o estado operacional",
            "as mensagens ficam ligadas ao pedido para coordenacao em contexto",
            "o perfil concentra configuracoes, conta e acessos auxiliares",
        ],
    )


def slide_client_gallery_2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Cliente",
        "Suporte e fluxo operacional",
        "Para alem da descoberta, o app cobre acompanhamento do pedido, suporte e resolucao de incidentes.",
    )
    add_phone_frame(
        slide,
        IMG_CLIENT_SUPPORT,
        Inches(0.85),
        Inches(1.7),
        Inches(4.2),
        Inches(4.75),
        "Suporte do cliente",
    )
    add_card(
        slide,
        Inches(5.35),
        Inches(1.8),
        Inches(6.85),
        Inches(1.72),
        "O que este bloco acrescenta",
        body="O cliente nao fica limitado a criar pedido e esperar. Existe canal formal para ajuda, reclamacao, duvida e mediação operacional.",
    )
    add_card(
        slide,
        Inches(5.35),
        Inches(3.82),
        Inches(6.85),
        Inches(2.1),
        "Fluxo resumido do cliente",
        bullets=[
            "entra no app, escolhe modo e procura servico",
            "abre pedido com contexto, localizacao e urgencia",
            "acompanha resposta, proposta, aceite e andamento",
            "usa mensagens, suporte, avaliacao e historico para fechar o ciclo",
        ],
    )


def slide_lifecycle():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Pedido",
        "Fluxo ponta a ponta dentro da plataforma",
        "O dominio do app cobre criacao, matching, proposta, aceite, execucao, confirmacao de valor, conclusao, cancelamento e no-show.",
    )
    add_flow_bar(
        slide,
        Inches(0.68),
        Inches(1.95),
        [
            ("Criacao", PRIMARY),
            ("Matching", NAVY),
            ("Proposta", PRIMARY),
            ("Aceite", NAVY),
            ("Execucao", PRIMARY),
            ("Conclusao", NAVY),
        ],
    )
    add_card(
        slide,
        Inches(0.70),
        Inches(3.0),
        Inches(3.8),
        Inches(2.3),
        "Inicio do pedido",
        bullets=[
            "cliente define titulo, descricao e categoria",
            "aplica urgencia, agenda ou pedido de orcamento",
            "localizacao pode ser automatica ou manual",
        ],
    )
    add_card(
        slide,
        Inches(4.78),
        Inches(3.0),
        Inches(3.8),
        Inches(2.3),
        "Coordenacao",
        bullets=[
            "prestador entra por matching ou convite",
            "pode aceitar e propor valor/faixa de preco",
            "cliente confirma e segue para execucao",
        ],
    )
    add_card(
        slide,
        Inches(8.86),
        Inches(3.0),
        Inches(3.8),
        Inches(2.3),
        "Fecho operacional",
        bullets=[
            "mensagens e mapa acompanham o trabalho",
            "historico e estados permitem auditoria",
            "cancelamento, reembolso e no-show sao tratados na plataforma",
        ],
    )


def slide_messages():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Mensagens",
        "Conversa real dentro do pedido",
        "As capturas abaixo mostram o fio da conversa aberto, do lado do cliente e do lado do prestador, ja dentro do ambiente de mensagens.",
    )
    add_phone_frame(
        slide,
        IMG_CLIENT_CHAT_THREAD,
        Inches(0.78),
        Inches(1.8),
        Inches(2.65),
        Inches(4.6),
        "Chat do cliente",
        caption_color=AMBER,
    )
    add_phone_frame(
        slide,
        IMG_PROVIDER_CHAT_THREAD,
        Inches(3.70),
        Inches(1.8),
        Inches(2.65),
        Inches(4.6),
        "Chat do prestador",
        caption_color=AMBER,
    )
    add_card(
        slide,
        Inches(6.72),
        Inches(1.82),
        Inches(5.45),
        Inches(2.0),
        "O que esta demonstracao evidencia",
        bullets=[
            "mensagens ficam contextualizadas no pedido",
            "cliente e prestador veem a mesma conversa por papeis distintos",
            "estrutura visual de bolhas, horario e cabecalho do pedido fica clara",
        ],
    )
    add_card(
        slide,
        Inches(6.72),
        Inches(4.02),
        Inches(5.45),
        Inches(2.1),
        "Exemplo usado nas capturas",
        bullets=[
            "cliente: Ola, tudo bem?",
            "prestador: Tudo e voce?",
            "objetivo: mostrar a area de conversa ja aberta e pronta para demo",
        ],
    )
    add_textbox(
        slide,
        Inches(0.80),
        Inches(6.55),
        Inches(5.7),
        Inches(0.28),
        "Nota: estas capturas foram geradas com o chat aberto para evitar qualquer slide com loading na area de mensagens.",
        font_size=9,
        color=MUTED,
    )


def slide_appendix_1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Outras telas",
        "Galeria complementar do modulo de mensagens e suporte",
        "Estas capturas entram no fim da apresentacao como reserva visual para demonstrar outras vistas do aplicativo.",
    )
    add_gallery_2x2(
        slide,
        [
            (IMG_CLIENT_MESSAGES, "Lista de mensagens do cliente"),
            (IMG_PROVIDER_MESSAGES, "Lista de mensagens do prestador"),
            (IMG_CLIENT_SUPPORT, "Suporte do cliente"),
            (IMG_PROVIDER_SUPPORT, "Suporte do prestador"),
        ],
    )


def slide_appendix_2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Outras telas",
        "Galeria complementar de navegacao e operacao",
        "Bloco final para fechar a demo com mais telas abertas do app e do backoffice.",
    )
    add_gallery_2x2(
        slide,
        [
            (IMG_ROLE, "Seletor de papel"),
            (IMG_CLIENT_ORDERS, "Pedidos do cliente"),
            (IMG_PROVIDER_PAYMENTS, "Pagamentos do prestador"),
            (IMG_ADMIN, "Painel administrativo"),
        ],
    )


def slide_provider_gallery_1():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Prestador",
        "Telas centrais da operacao do prestador",
        "O prestador controla disponibilidade, recebe trabalhos, acompanha conversas e gere a sua conta.",
    )
    add_gallery_2x2(
        slide,
        [
            (IMG_PROVIDER_HOME, "Home do prestador"),
            (IMG_PROVIDER_JOBS, "Trabalhos"),
            (IMG_PROVIDER_MESSAGES, "Mensagens"),
            (IMG_PROVIDER_PROFILE, "Perfil"),
        ],
    )
    add_card(
        slide,
        Inches(6.35),
        Inches(1.75),
        Inches(6.0),
        Inches(4.95),
        "Leitura do fluxo do prestador",
        bullets=[
            "liga o modo online e configura a area de atuacao",
            "ve oportunidades e trabalhos em curso",
            "coordena o servico pelas mensagens ligadas ao pedido",
            "usa o perfil para pagamentos, configuracoes e suporte",
        ],
    )


def slide_provider_gallery_2():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Prestador",
        "Pagamentos, onboarding e validacao",
        "Esta camada transforma o app num marketplace transacional e nao apenas num diretório de servicos.",
    )
    add_phone_frame(
        slide,
        IMG_PROVIDER_PAYMENTS,
        Inches(0.85),
        Inches(1.7),
        Inches(4.2),
        Inches(4.75),
        "Pagamentos do prestador",
    )
    add_card(
        slide,
        Inches(5.38),
        Inches(1.86),
        Inches(6.8),
        Inches(1.85),
        "Componentes deste modulo",
        bullets=[
            "Stripe Connect Express para onboarding financeiro",
            "assinaturas Basic e Pro",
            "upload documental e estado de KYC",
        ],
    )
    add_card(
        slide,
        Inches(5.38),
        Inches(4.02),
        Inches(6.8),
        Inches(2.0),
        "Impacto operacional",
        bullets=[
            "liberta pagamentos e profissionaliza o lado do prestador",
            "permite controlar elegibilidade e conformidade documental",
            "cria base para retencao via subscricao e recorrencia",
        ],
    )


def slide_trust():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Confianca",
        "Camadas que sustentam a experiencia completa",
        "A proposta do app depende de coordenacao, pagamento, localizacao, avaliacao e suporte a incidentes.",
    )
    add_card(
        slide,
        Inches(0.65),
        Inches(1.82),
        Inches(2.95),
        Inches(1.95),
        "Comunicacao",
        bullets=["chat por pedido", "nao lidas, typing e vistos", "anexos e localizacao"],
    )
    add_card(
        slide,
        Inches(3.84),
        Inches(1.82),
        Inches(2.95),
        Inches(1.95),
        "Localizacao",
        bullets=["morada do pedido", "tracking do prestador", "apoio ao encontro fisico"],
    )
    add_card(
        slide,
        Inches(7.03),
        Inches(1.82),
        Inches(2.75),
        Inches(1.95),
        "Pagamento",
        bullets=["Payment Sheet no cliente", "Connect no prestador", "subscricao e monetizacao"],
    )
    add_card(
        slide,
        Inches(10.02),
        Inches(1.82),
        Inches(2.63),
        Inches(1.95),
        "Controlo",
        bullets=["KYC", "no-show", "reembolso", "historico"],
    )
    add_card(
        slide,
        Inches(0.65),
        Inches(4.12),
        Inches(6.0),
        Inches(1.85),
        "Do ponto de vista do cliente",
        bullets=[
            "segue o pedido com mais previsibilidade",
            "tem evidencia historica do que foi combinado",
            "recorre a suporte e avaliacao quando necessario",
        ],
    )
    add_card(
        slide,
        Inches(6.9),
        Inches(4.12),
        Inches(5.75),
        Inches(1.85),
        "Do ponto de vista da plataforma",
        bullets=[
            "controla risco operacional e qualidade de servico",
            "tem mecanismos de monetizacao e elegibilidade",
            "ganha visibilidade sobre performance e problemas",
        ],
    )


def slide_admin():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Backoffice",
        "Tela administrativa e controlo da operacao",
        "A plataforma tambem tem camada de administracao para acompanhar metricas, suporte e gestao de excecoes.",
    )
    add_phone_frame(
        slide,
        IMG_ADMIN,
        Inches(0.85),
        Inches(1.7),
        Inches(4.25),
        Inches(4.75),
        "Painel administrativo",
    )
    add_card(
        slide,
        Inches(5.45),
        Inches(1.84),
        Inches(6.75),
        Inches(1.9),
        "O que aparece no painel",
        bullets=[
            "metricas e indicadores operacionais",
            "suporte, no-show e moderacao",
            "visao central da plataforma",
        ],
    )
    add_card(
        slide,
        Inches(5.45),
        Inches(4.04),
        Inches(6.75),
        Inches(2.0),
        "Valor desta camada",
        bullets=[
            "permite mediar excecoes e conflitos",
            "apoia a qualidade de servico e a confianca do marketplace",
            "fecha o ciclo entre cliente, prestador e operacao",
        ],
    )


def slide_summary():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_header(
        slide,
        "Resumo",
        "Leitura final do aplicativo",
        "O ChegaJa foi desenhado como um fluxo operacional completo, nao apenas como vitrine de prestadores.",
    )
    add_card(
        slide,
        Inches(0.75),
        Inches(1.95),
        Inches(3.85),
        Inches(2.15),
        "Cliente",
        bullets=[
            "descobre servico",
            "abre e acompanha pedido",
            "fala com o prestador e pede suporte",
        ],
        fill=SOFT,
    )
    add_card(
        slide,
        Inches(4.75),
        Inches(1.95),
        Inches(3.85),
        Inches(2.15),
        "Prestador",
        bullets=[
            "recebe oportunidade relevante",
            "opera trabalhos e mensagens",
            "gere pagamentos e documentacao",
        ],
        fill=SOFT,
    )
    add_card(
        slide,
        Inches(8.75),
        Inches(1.95),
        Inches(3.85),
        Inches(2.15),
        "Plataforma",
        bullets=[
            "orquestra matching e estados do pedido",
            "monetiza via pagamento e subscricao",
            "controla suporte, risco e conformidade",
        ],
        fill=SOFT,
    )
    add_card(
        slide,
        Inches(0.75),
        Inches(4.55),
        Inches(11.85),
        Inches(1.35),
        "Resultado",
        body="Esta versao da apresentacao passa a incluir as capturas reais disponiveis do app para mostrar o fluxo visual do produto de forma muito mais completa.",
    )


OUT_DIR.mkdir(parents=True, exist_ok=True)

slide_cover()
slide_overview()
slide_entry()
slide_client_gallery_1()
slide_client_gallery_2()
slide_lifecycle()
slide_messages()
slide_provider_gallery_1()
slide_provider_gallery_2()
slide_trust()
slide_admin()
slide_summary()
slide_appendix_1()
slide_appendix_2()

prs.save(OUT_FILE_FIXED)
print(OUT_FILE_FIXED)

for legacy_path in (OUT_FILE, OUT_FILE_FULL):
    try:
        prs.save(legacy_path)
        print(legacy_path)
    except PermissionError:
        print(f"LOCKED: {legacy_path}")
